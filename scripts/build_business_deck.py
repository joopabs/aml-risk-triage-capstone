"""Build reports/slides/business_deck.pptx from reports/slides/business_deck_outline.md (python-pptx).

Outline grammar: "## Slide N — Title" starts a slide; "- " bullets; "| a | b |" table rows;
"> " blockquote (disclaimer, rendered as a footer box); "*...*" single-line italics become speaker notes;
a line starting with "**Label:**" becomes a bold lead-in. The first slide is the title slide.
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
RED = RGBColor(0xDD, 0x51, 0x43)
GREY = RGBColor(0x55, 0x55, 0x55)


def parse(md: str) -> list[dict]:
    slides, cur = [], None
    for line in md.splitlines():
        m = re.match(r"^## Slide (\d+) — (.*)$", line)
        if m:
            cur = {"title": m.group(2).strip(), "lines": [], "table": [], "quote": [], "notes": []}
            slides.append(cur)
            continue
        if cur is None or not line.strip() or line.strip() == "---":
            continue
        if line.startswith("|"):
            if set(line.replace("|", "").strip()) <= {"-", " "}:
                continue
            cur["table"].append([c.strip() for c in line.strip().strip("|").split("|")])
        elif line.startswith("> "):
            cur["quote"].append(line[2:].strip())
        elif re.match(r"^\*[^*].*\*$", line.strip()):
            cur["notes"].append(line.strip().strip("*"))
        else:
            cur["lines"].append(line.strip())
    return slides


def clean(text: str) -> str:
    return re.sub(r"[*`]", "", text)


def add_text(tf, lines: list[str], size: int = 16):
    first = True
    for raw in lines:
        bullet = raw.startswith("- ") or bool(re.match(r"^\d+\. ", raw))
        text = clean(re.sub(r"^(- |\d+\. )", "", raw))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = ("• " if bullet else "") + text
        p.font.size = Pt(size)
        if raw.startswith("**") and ":**" in raw:
            p.font.bold = True
        p.space_after = Pt(6)


def build(outline: Path, out: Path) -> int:
    slides = parse(outline.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    for i, s in enumerate(slides):
        sl = prs.slides.add_slide(blank)
        # title
        tb = sl.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0)).text_frame
        tb.text = clean(s["title"]) if i else clean(s["lines"][0]) if s["lines"] else s["title"]
        tb.paragraphs[0].font.size = Pt(34 if i == 0 else 28)
        tb.paragraphs[0].font.bold = True
        tb.paragraphs[0].font.color.rgb = NAVY
        body_lines = s["lines"][1:] if i == 0 else s["lines"]
        top = Inches(1.5)
        if body_lines:
            h = Inches(2.6 if s["table"] else 4.6)
            body = sl.shapes.add_textbox(Inches(0.7), top, Inches(11.9), h).text_frame
            body.word_wrap = True
            add_text(body, body_lines, size=20 if i == 0 else 16)
            top = top + h + Inches(0.1)
        if s["table"]:
            rows, cols = len(s["table"]), len(s["table"][0])
            height = Inches(0.42) * rows
            shape = sl.shapes.add_table(rows, cols, Inches(0.7), top, Inches(11.9), height)
            for r, row in enumerate(s["table"]):
                for c, val in enumerate(row):
                    cell = shape.table.cell(r, c)
                    cell.text = clean(val)
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(13)
                        p.font.bold = r == 0 or val.startswith("**")
            top = top + height + Inches(0.15)
        if s["quote"]:
            q = sl.shapes.add_textbox(
                Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.9)
            ).text_frame
            q.word_wrap = True
            q.text = clean(" ".join(s["quote"]))
            q.paragraphs[0].font.size = Pt(11)
            q.paragraphs[0].font.italic = True
            q.paragraphs[0].font.color.rgb = RED
        foot = sl.shapes.add_textbox(
            Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.35)
        ).text_frame
        foot.text = f"Synthetic PaySim data · educational decision-support prototype · no AML determination · slide {i + 1}/{len(slides)}"
        foot.paragraphs[0].font.size = Pt(9)
        foot.paragraphs[0].font.color.rgb = GREY
        if s["notes"]:
            sl.notes_slide.notes_text_frame.text = clean(" ".join(s["notes"]))
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"wrote {out} ({len(slides)} slides)")
    return len(slides)


if __name__ == "__main__":
    src = (
        Path(sys.argv[1]) if len(sys.argv) > 1 else Path("reports/slides/business_deck_outline.md")
    )
    dst = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("reports/slides/business_deck.pptx")
    build(src, dst)
